from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ── 颜色（完全复刻参考PPT）
C_RED    = RGBColor(0xA5, 0x0A, 0x0D)  # 参考PPT线条色
C_RED2   = RGBColor(0xC0, 0x00, 0x00)  # 标签填充色
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK   = RGBColor(0x1A, 0x1A, 0x1A)
C_GRAY   = RGBColor(0x55, 0x55, 0x55)
C_LGRAY  = RGBColor(0xE8, 0xE8, 0xE8)
C_CARD   = RGBColor(0xFF, 0xFF, 0xFF)

BG   = r"C:\Netease_Ai\03_document_processing\ppt\outputs\assets\bg_master.jpg"
CHARTS = r"C:\Netease_Ai\03_document_processing\ppt\outputs\charts"

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

def blank():
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(slide):
    """铺全幅背景图"""
    slide.shapes.add_picture(BG, 0, 0, W, H)

def rect(slide, l, t, w, h, fill, lc=None, lw=Pt(0)):
    sh = slide.shapes.add_shape(1, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if lc: sh.line.color.rgb = lc; sh.line.width = lw
    else:  sh.line.fill.background()
    return sh

def tb(slide, text, l, t, w, h, sz=16, bold=False,
       color=C_WHITE, align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(sz); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = color
    r.font.name = "微软雅黑"
    return box

def ml(slide, lines, l, t, w, h, sz=13, color=C_DARK, sp=76200):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Emu(sp)
        r = p.add_run(); r.text = line
        r.font.size = Pt(sz); r.font.name = "微软雅黑"
        r.font.color.rgb = color
    return box

def tag(slide, text, l, t, w=Inches(1.38), h=Inches(0.33)):
    """红色标签（居中，复刻原版）"""
    rect(slide, l, t, w, h, C_RED2)
    tb(slide, text, l, t+Inches(0.04), w, h,
       sz=12, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

def hline(slide, l, t, w):
    """水平红线（复刻原版 line_width=25400）"""
    ln = slide.shapes.add_connector(1, l, t, l+w, t)
    ln.line.color.rgb = C_RED
    ln.line.width = Emu(25400)

def title_bar(slide, title, sub=None):
    """页面标题区（白色文字浮于红色背景）"""
    tb(slide, title,
       Inches(0.494), Inches(0.328), W-Inches(1), Inches(0.52),
       sz=24, bold=True, color=C_WHITE)
    hline(slide, Inches(0.583), Inches(0.966), Inches(0.583))
    if sub:
        tb(slide, sub,
           Inches(0.494), Inches(1.02), W-Inches(1), Inches(0.35),
           sz=13, color=C_WHITE)

def white_card(slide, l, t, w, h):
    """白色大卡片浮于红底"""
    sh = slide.shapes.add_shape(1, l, t, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = C_CARD
    sh.line.fill.background()
    return sh

def card_with_tag(slide, l, t, w, h, tag_txt, title, desc, sz_title=15, sz_desc=12):
    """完整卡片：白底 + 红色标签 + 标题 + 描述"""
    white_card(slide, l, t, w, h)
    tag(slide, tag_txt, l+Inches(0.15), t+Inches(0.15), Inches(1.2), Inches(0.31))
    tb(slide, title,
       l+Inches(0.15), t+Inches(0.55), w-Inches(0.3), Inches(0.42),
       sz=sz_title, bold=True, color=C_DARK)
    tb(slide, desc,
       l+Inches(0.15), t+Inches(1.0), w-Inches(0.3), h-Inches(1.15),
       sz=sz_desc, color=C_GRAY)

def chart_card(slide, fname, l, t, w, h, cap):
    """图表卡片：白底 + 图表 + 说明"""
    white_card(slide, l, t, w, h)
    slide.shapes.add_picture(
        os.path.join(CHARTS, fname),
        l+Inches(0.1), t+Inches(0.1),
        w-Inches(0.2), h-Inches(0.55))
    tb(slide, cap,
       l, t+h-Inches(0.42), w, Inches(0.35),
       sz=9, italic=True, color=C_GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════
# P01  封面
# ═══════════════════════════════════════
s1 = blank(); bg(s1)
tb(s1, "AI 技术在游戏交互环节的运用",
   Inches(0.583), Inches(2.3), Inches(10), Inches(1.2),
   sz=40, bold=True, color=C_WHITE)
hline(s1, Inches(0.583), Inches(3.6), Inches(1.2))
tb(s1, "现状 · 案例 · 启发",
   Inches(0.583), Inches(3.75), Inches(8), Inches(0.5),
   sz=20, color=C_WHITE)
tb(s1, "分享人：XXX　　2026.03",
   Inches(0.583), Inches(4.45), Inches(6), Inches(0.4),
   sz=14, color=C_WHITE)


# ═══════════════════════════════════════
# P02  为什么谈这个 + 柱状图
# ═══════════════════════════════════════
s2 = blank(); bg(s2)
title_bar(s2, "为什么现在谈这个？",
          "AI 正在改变游戏交互的底层逻辑——数据说话")

card_with_tag(s2, Inches(0.583), Inches(1.35),
              Inches(3.05), Inches(5.75),
              "过去", "规则驱动",
              "· 设计师写死规则\n· NPC 台词固定\n· 难度上线就定死\n· 新手教程人人一样",
              sz_title=16, sz_desc=13)

tb(s2, "→",
   Inches(3.75), Inches(3.8), Inches(0.5), Inches(0.6),
   sz=28, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

card_with_tag(s2, Inches(4.35), Inches(1.35),
              Inches(3.05), Inches(5.75),
              "现在", "玩家驱动",
              "· AI 实时感知行为\n· NPC 实时生成对话\n· 难度随水平浮动\n· 引导路径因人而异",
              sz_title=16, sz_desc=13)

chart_card(s2, "chart_p02_compare.png",
           Inches(7.6), Inches(1.35),
           Inches(5.3), Inches(5.75),
           "数据来源：行业研究报告综合均值（仅供参考）")


# ═══════════════════════════════════════
# P03  讨论范围
# ═══════════════════════════════════════
s3 = blank(); bg(s3)
title_bar(s3, "本次讨论的范围", "我们说的「游戏交互」包括哪些？")

topics = [
    ("对话互动", "NPC 对话与互动",    "AI 生成对话，告别写死台词"),
    ("新手引导", "新手引导与教学",    "个性化引导路径，因人而异"),
    ("UI反馈",  "UI 反馈与操作体验", "预判意图，提前响应玩家操作"),
    ("难度节奏", "难度曲线与节奏",    "DDA 动态调节，保持心流体验"),
    ("情绪感知", "玩家情绪感知",      "读懂情绪，调整游戏状态（前沿）"),
    ("不涉及",   "内容生成/数据分析", "本次暂不覆盖"),
]
bw = Inches(3.95); bh = Inches(2.3)
gx = Inches(0.22); gy = Inches(0.2)
ox = Inches(0.583); oy = Inches(1.35)
for i, (t, ti, d) in enumerate(topics):
    col = i % 3; row = i // 3
    card_with_tag(s3,
                  ox + col*(bw+gx), oy + row*(bh+gy),
                  bw, bh, t, ti, d)


# ═══════════════════════════════════════
# P04  5大方向 + 雷达图
# ═══════════════════════════════════════
s4 = blank(); bg(s4)
title_bar(s4, "5 大应用方向总览", "AI × 游戏交互——行业成熟度与落地情况")

dirs = [
    ("方向①", "智能 NPC 对话",  "让NPC真正能聊天\n成熟度：72%"),
    ("方向②", "动态难度调节",   "感知水平自动调节\n成熟度：85%"),
    ("方向③", "个性化引导",     "新手教学因人而异\n成熟度：68%"),
    ("方向④", "行为预测与反馈", "提前知道你下一步\n成熟度：60%"),
    ("方向⑤", "情感感知交互",   "读懂情绪调节节奏\n成熟度：28%"),
]
bw = Inches(2.3); bh = Inches(3.7); g = Inches(0.17)
for i, (t, ti, d) in enumerate(dirs):
    card_with_tag(s4, Inches(0.583)+i*(bw+g), Inches(1.35),
                  bw, bh, t, ti, d)

chart_card(s4, "chart_p04_radar.png",
           Inches(3.2), Inches(5.2),
           Inches(6.9), Inches(2.15),
           "各方向行业落地成熟度雷达图（基于公开数据综合评估）")


# ═══════════════════════════════════════
# P05  NPC + DDA
# ═══════════════════════════════════════
s5 = blank(); bg(s5)
title_bar(s5, "方向①② — 智能 NPC + 动态难度",
          "让 NPC 会说话，让游戏懂你的水平")

# 左侧大白卡
white_card(s5, Inches(0.583), Inches(1.35), Inches(5.95), Inches(5.8))
tag(s5, "方向①", Inches(0.73), Inches(1.5), Inches(1.1), Inches(0.31))
tb(s5, "智能 NPC 对话",
   Inches(0.73), Inches(1.9), Inches(5.6), Inches(0.44),
   sz=18, bold=True, color=C_DARK)
rect(s5, Inches(0.73), Inches(2.4), Inches(5.6), Pt(1), C_LGRAY)
ml(s5, [
    "技  术：LLM 驱动，NPC 实时生成对话，告别写死剧本",
    "案  例：Cyberpunk 2077 GPT-4 实验",
    "           Inworld AI 引擎已被多款游戏接入",
    "设计影响：剧情分支从【树状】变为【开放式】",
    "           叙事边界扩展，一致性把控是新挑战",
], Inches(0.73), Inches(2.5), Inches(5.6), Inches(2.8), sz=14)

tag(s5, "方向②", Inches(0.73), Inches(5.35), Inches(1.1), Inches(0.31))
tb(s5, "数据支撑",
   Inches(2.0), Inches(5.38), Inches(4.2), Inches(0.3),
   sz=12, bold=True, color=C_RED2)
tb(s5, "行业成熟度 72%  |  Inworld 已接入 50+ 款游戏",
   Inches(0.73), Inches(5.75), Inches(5.6), Inches(0.3),
   sz=12, color=C_GRAY)

# 右侧大白卡
white_card(s5, Inches(6.8), Inches(1.35), Inches(5.95), Inches(5.8))
tag(s5, "方向②", Inches(6.94), Inches(1.5), Inches(1.1), Inches(0.31))
tb(s5, "动态难度调节（DDA）",
   Inches(6.94), Inches(1.9), Inches(5.6), Inches(0.44),
   sz=18, bold=True, color=C_DARK)
rect(s5, Inches(6.94), Inches(2.4), Inches(5.6), Pt(1), C_LGRAY)
ml(s5, [
    "技  术：实时监测死亡率/操作精度，自动微调敌人强度",
    "案  例：《生化危机》系列 —— DDA 已用 20 年",
    "           《FIFA》对手 AI 随玩家水平实时浮动",
    "设计影响：难度调参不再全靠策划手动完成",
    "           AI 接管「让玩家保持心流」这件事",
], Inches(6.94), Inches(2.5), Inches(5.6), Inches(2.8), sz=14)

tag(s5, "数据", Inches(6.94), Inches(5.35), Inches(0.8), Inches(0.31))
tb(s5, "行业成熟度 85%  |  已有 60%+ 的 AAA 游戏采用",
   Inches(7.85), Inches(5.38), Inches(4.7), Inches(0.3),
   sz=12, color=C_GRAY)


# ═══════════════════════════════════════
# P06  个性化引导 + AB折线图
# ═══════════════════════════════════════
s6 = blank(); bg(s6)
title_bar(s6, "方向③④ — 个性化引导 + 行为预测",
          "新手不再走同一条路，系统提前知道你要做什么")

white_card(s6, Inches(0.583), Inches(1.35), Inches(5.95), Inches(5.8))
tag(s6, "方向③", Inches(0.73), Inches(1.5), Inches(1.1), Inches(0.31))
tb(s6, "个性化新手引导",
   Inches(0.73), Inches(1.9), Inches(5.6), Inches(0.44),
   sz=18, bold=True, color=C_DARK)
rect(s6, Inches(0.73), Inches(2.4), Inches(5.6), Pt(1), C_LGRAY)
ml(s6, [
    "技  术：行为聚类识别玩家类型，推送匹配引导路径",
    "案  例：《原神》多路径引导持续迭代，留存显著提升",
    "           King 把多路径 AB 测试做成标配工具",
    "设计影响：从「一套走天下」变成「多条路径并行」",
], Inches(0.73), Inches(2.5), Inches(5.6), Inches(1.9), sz=14)

# AB折线图嵌入左卡底部
slide_l = Inches(0.73); slide_t = Inches(4.45)
slide_w = Inches(5.6);  slide_h = Inches(2.4)
s6.shapes.add_picture(
    os.path.join(CHARTS, "chart_p06_ab.png"),
    slide_l, slide_t, slide_w, slide_h)

white_card(s6, Inches(6.8), Inches(1.35), Inches(5.95), Inches(5.8))
tag(s6, "方向④", Inches(6.94), Inches(1.5), Inches(1.1), Inches(0.31))
tb(s6, "行为预测与反馈",
   Inches(6.94), Inches(1.9), Inches(5.6), Inches(0.44),
   sz=18, bold=True, color=C_DARK)
rect(s6, Inches(6.94), Inches(2.4), Inches(5.6), Pt(1), C_LGRAY)
ml(s6, [
    "技  术：分析操作序列预测下一步意图，提前触发提示",
    "案  例：腾讯 AI Lab — 预测流失节点，提前推送挽留",
    "           《英雄联盟》— 实时走位分析优化提示时机",
    "设计影响：UI 反馈从「响应操作」变成「预判意图」",
    "           用户体验更流畅，操作摩擦大幅降低",
], Inches(6.94), Inches(2.5), Inches(5.6), Inches(3.8), sz=14)


# ═══════════════════════════════════════
# P07  情感感知交互
# ═══════════════════════════════════════
s7 = blank(); bg(s7)
title_bar(s7, "方向⑤ — 情感感知交互（前沿方向）",
          "下一步：游戏能「感知」你的情绪")

tech = [
    ("摄像头",   "面部表情识别\nAffectiva / Emotion AI"),
    ("语音",     "声纹分析\n识别焦虑/兴奋/沮丧"),
    ("生理信号", "心率/皮肤电\n穿戴设备/XR场景"),
]
for i, (t, d) in enumerate(tech):
    card_with_tag(s7,
                  Inches(0.583)+i*Inches(2.75), Inches(1.35),
                  Inches(2.5), Inches(2.1), t, d, "")

tb(s7, "→  情感模型  →  游戏动态响应",
   Inches(9.0), Inches(1.9), Inches(4.0), Inches(0.6),
   sz=15, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

# 下方两列
white_card(s7, Inches(0.583), Inches(3.7), Inches(5.95), Inches(3.45))
tb(s7, "已有实验",
   Inches(0.73), Inches(3.85), Inches(5.6), Inches(0.4),
   sz=16, bold=True, color=C_DARK)
rect(s7, Inches(0.73), Inches(4.32), Inches(5.6), Pt(1), C_LGRAY)
ml(s7, [
    "· EA 研究院：用情感识别优化恐怖游戏惊吓时机",
    "· VR《生化危机 7》：测试版接入心率监测",
    "· Xbox 无障碍方向：情感感知辅助残障玩家",
], Inches(0.73), Inches(4.42), Inches(5.6), Inches(2.4), sz=14)

white_card(s7, Inches(6.8), Inches(3.7), Inches(5.95), Inches(3.45))
tb(s7, "现实限制",
   Inches(6.94), Inches(3.85), Inches(5.6), Inches(0.4),
   sz=16, bold=True, color=C_DARK)
rect(s7, Inches(6.94), Inches(4.32), Inches(5.6), Pt(1), C_LGRAY)
ml(s7, [
    "· 隐私合规门槛高",
    "· 移动端硬件暂不具备规模化",
    "· 复杂环境识别准确率仍不稳定",
    "· 现在不一定要做，但值得在设计框架里预留接口",
], Inches(6.94), Inches(4.42), Inches(5.6), Inches(2.4), sz=14)


# ═══════════════════════════════════════
# P08  启发 + 饼图
# ═══════════════════════════════════════
s8 = blank(); bg(s8)
title_bar(s8, "对我们的启发", "作为设计师/策划，我们现在能做什么？")

cols = [
    ("可以交给AI",  ["难度曲线的微调参数", "新手引导路径分支测试", "玩家行为数据分析与分类"]),
    ("需要人来把关", ["交互体验的情感基调", "NPC 对话的世界观一致性", "最终设计方向的价值判断"]),
    ("近期可以尝试", ["① Inworld AI 做 NPC 对话原型", "② 预设 2 条引导路径做 AB 测", "③ 关注竞品引导/难度版本迭代"]),
]
cw = Inches(3.0); cg = Inches(0.22); sy = Inches(1.35)
for i, (ct, items) in enumerate(cols):
    lx = Inches(0.583)+i*(cw+cg)
    white_card(s8, lx, sy, cw, Inches(5.8))
    rect(s8, lx, sy, cw, Inches(0.42), C_RED2)
    tb(s8, ct, lx, sy+Inches(0.05), cw, Inches(0.37),
       sz=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    for j, it in enumerate(items):
        ty = sy+Inches(0.55)+j*Inches(1.55)
        rect(s8, lx+Inches(0.15), ty, Inches(0.06), Inches(1.35), C_RED2)
        tb(s8, it, lx+Inches(0.3), ty+Inches(0.35),
           cw-Inches(0.45), Inches(0.9), sz=13, color=C_DARK)

chart_card(s8, "chart_p08_pie.png",
           Inches(10.0), Inches(1.35),
           Inches(2.95), Inches(5.8),
           "引入AI后交互设计工作量结构")


# ═══════════════════════════════════════
# P09  总结 + 趋势图
# ═══════════════════════════════════════
s9 = blank(); bg(s9)
title_bar(s9, "核心结论", "三句话总结 + 行业趋势")

conclusions = [
    ("01", "AI 在游戏交互中已不是「未来技术」",
           "NPC 对话、难度调节等能力今天就能用"),
    ("02", "设计师/策划的角色在升级",
           "从「规则制定者」变成「AI 行为的把关人」"),
    ("03", "现在最值得关注的方向",
           "智能 NPC 对话  +  个性化引导"),
]
for i, (num, ti, sub) in enumerate(conclusions):
    ty = Inches(1.35)+i*Inches(1.5)
    white_card(s9, Inches(0.583), ty, Inches(6.3), Inches(1.35))
    rect(s9, Inches(0.583), ty, Inches(0.55), Inches(1.35), C_RED2)
    tb(s9, num,
       Inches(0.583), ty+Inches(0.45), Inches(0.55), Inches(0.45),
       sz=16, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    tb(s9, ti,
       Inches(1.25), ty+Inches(0.15), Inches(5.5), Inches(0.45),
       sz=16, bold=True, color=C_DARK)
    tb(s9, sub,
       Inches(1.25), ty+Inches(0.65), Inches(5.5), Inches(0.5),
       sz=13, color=C_GRAY)

chart_card(s9, "chart_p09_trend.png",
           Inches(7.2), Inches(1.25),
           Inches(5.75), Inches(5.85),
           "数据来源：Newzoo / IDC 游戏AI报告综合整理")

hline(s9, Inches(0.583), Inches(6.55), W-Inches(1.2))
tb(s9, "你们项目里，有没有哪个交互环节特别痛？NPC 太死板 / 引导流失 / 手动调参——这些正是 AI 最先能帮上的地方。",
   Inches(0.583), Inches(6.7), Inches(6.5), Inches(0.7),
   sz=12, italic=True, color=C_WHITE)


# ═══════════════════════════════════════
# 保存
# ═══════════════════════════════════════
OUT = r"C:\Netease_Ai\03_document_processing\ppt\outputs\20260326_AI技术在游戏交互环节的运用_v2.0.pptx"
prs.save(OUT)
print("Saved: " + OUT)