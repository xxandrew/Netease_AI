from pptx import Presentation
from pptx.util import Pt

prs = Presentation(r"C:/Users/xxn4472/Desktop/test1/MC_UXD-AI0323.pptx")
lines = []
lines.append("SIZE=%.2fx%.2f PAGES=%d" % (prs.slide_width.inches, prs.slide_height.inches, len(prs.slides)))

def info(sh, d=0):
    pad = "  " * d
    l = round((sh.left or 0)/914400, 2)
    t = round((sh.top or 0)/914400, 2)
    w = round((sh.width or 0)/914400, 2)
    h = round((sh.height or 0)/914400, 2)
    lines.append("%sSH[%s] %s (%s,%s) %sx%s" % (pad, sh.shape_type, sh.name, l, t, w, h))
    try:
        fi = sh.fill
        lines.append("%s fill_t=%s" % (pad, fi.type))
        try:
            lines.append("%s fill_rgb=%s" % (pad, fi.fore_color.rgb))
        except:
            pass
    except:
        pass
    if sh.has_text_frame:
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                if r.text.strip():
                    try:
                        c = str(r.font.color.rgb)
                    except:
                        c = "inh"
                    lines.append("%s T=%r sz=%s b=%s c=%s fn=%s" % (pad, r.text[:50], r.font.size, r.font.bold, c, r.font.name))
    if sh.shape_type == 6:
        for ch in sh.shapes:
            info(ch, d+1)

for i, sl in enumerate(prs.slides):
    lines.append("=== Slide %d ===" % (i+1))
    bg = sl.background.fill
    lines.append("bg_type=%s" % bg.type)
    try:
        lines.append("bg_rgb=%s" % bg.fore_color.rgb)
    except:
        lines.append("bg_rgb=NA")
    for sh in sl.shapes:
        info(sh)

with open("ref_out.txt", "w", encoding="utf-8") as f:
    f.write("\n".join(lines))
print("DONE")
