from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import json, os

prs = Presentation(r"C:/Users/xxn4472/Desktop/test1/MC_UXD-AI0323.pptx")

out = []
out.append("=== 基础信息 ===")
out.append("幻灯片尺寸: %.4f x %.4f emu = %.2f x %.2f inches" % (
    prs.slide_width, prs.slide_height,
    prs.slide_width/914400, prs.slide_height/914400))
out.append("页数: %d" % len(prs.slides))

# 读取 slide master 主题色
master = prs.slide_masters[0]
out.append("\n=== Slide Master 背景 ===")
bg = master.background.fill
out.append("master bg type: %s" % bg.type)
try: out.append("master bg rgb: %s" % bg.fore_color.rgb)
except: out.append("master bg rgb: N/A")

# 读取 master 中所有形状颜色
out.append("\n=== Master Shapes ===")
for sh in master.shapes:
    l=round((sh.left or 0)/914400,3); t=round((sh.top or 0)/914400,3)
    w=round((sh.width or 0)/914400,3); h=round((sh.height or 0)/914400,3)
    line = "  [%s] %s (%s,%s) %sx%s" % (sh.shape_type, sh.name, l, t, w, h)
    try:
        fi=sh.fill
        line += " fill=%s" % fi.type
        try: line += " rgb=%s" % fi.fore_color.rgb
        except: pass
    except: pass
    out.append(line)
    if sh.has_text_frame:
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                if r.text.strip():
                    try: c = str(r.font.color.rgb)
                    except: c = "inh"
                    out.append("    T=%r sz=%s b=%s c=%s fn=%s" % (
                        r.text[:50], r.font.size, r.font.bold, c, r.font.name))

# 读取每页 layout 颜色
out.append("\n=== Layouts ===")
for layout in master.slide_layouts:
    bg2 = layout.background.fill
    line = "Layout: %s bg_type=%s" % (layout.name, bg2.type)
    try: line += " bg_rgb=%s" % bg2.fore_color.rgb
    except: pass
    out.append(line)
    for sh in layout.shapes:
        l=round((sh.left or 0)/914400,3); t=round((sh.top or 0)/914400,3)
        w=round((sh.width or 0)/914400,3); h=round((sh.height or 0)/914400,3)
        sline = "  [%s] %s (%s,%s) %sx%s" % (sh.shape_type, sh.name, l, t, w, h)
        try:
            fi=sh.fill
            sline += " fill=%s" % fi.type
            try: sline += " rgb=%s" % fi.fore_color.rgb
            except: pass
        except: pass
        if sh.has_text_frame:
            for p in sh.text_frame.paragraphs:
                for r in p.runs:
                    if r.text.strip():
                        try: c = str(r.font.color.rgb)
                        except: c = "inh"
                        sline += " | T=%r sz=%s fn=%s c=%s" % (r.text[:30], r.font.size, r.font.name, c)
        out.append(sline)

# 每页详细信息
out.append("\n=== 每页详细 ===")
def dump_shape(sh, depth=0):
    pad = "  " * depth
    l=round((sh.left or 0)/914400,3); t=round((sh.top or 0)/914400,3)
    w=round((sh.width or 0)/914400,3); h=round((sh.height or 0)/914400,3)
    lines = []
    lines.append("%s[%s] %s | pos=(%s,%s) size=%sx%s" % (pad, sh.shape_type, sh.name, l, t, w, h))
    try:
        fi = sh.fill
        lines.append("%s  fill_type=%s" % (pad, fi.type))
        try: lines.append("%s  fill_rgb=%s" % (pad, fi.fore_color.rgb))
        except: pass
        try: lines.append("%s  fill_theme=%s" % (pad, fi.theme_color))
        except: pass
    except: pass
    # 线条
    try:
        ln = sh.line
        lines.append("%s  line_width=%s" % (pad, ln.width))
        try: lines.append("%s  line_rgb=%s" % (pad, ln.color.rgb))
        except: pass
    except: pass
    # 文本
    if sh.has_text_frame:
        tf = sh.text_frame
        lines.append("%s  auto_size=%s margin_l=%s margin_t=%s" % (
            pad, tf.auto_size, tf.margin_left, tf.margin_top))
        for pi, p in enumerate(tf.paragraphs):
            pline = "%s  PARA[%d] align=%s space_before=%s space_after=%s" % (
                pad, pi, p.alignment, p.space_before, p.space_after)
            lines.append(pline)
            for ri, r in enumerate(p.runs):
                if r.text.strip():
                    try: c = str(r.font.color.rgb)
                    except: c = "inh"
                    lines.append("%s    RUN[%d] T=%r sz=%s b=%s i=%s c=%s fn=%s" % (
                        pad, ri, r.text[:60], r.font.size, r.font.bold,
                        r.font.italic, c, r.font.name))
    # 子shape
    if sh.shape_type == 6:
        for ch in sh.shapes:
            lines.extend(dump_shape(ch, depth+1))
    return lines

for idx, slide in enumerate(prs.slides):
    out.append("\n--- Slide %d (layout=%s) ---" % (idx+1, slide.slide_layout.name))
    bg = slide.background.fill
    out.append("bg_type=%s" % bg.type)
    try: out.append("bg_rgb=%s" % bg.fore_color.rgb)
    except: out.append("bg_rgb=inherited")
    for sh in slide.shapes:
        out.extend(dump_shape(sh))

with open(r"C:\Netease_Ai\ref_deep.txt", "w", encoding="utf-8") as f:
    f.write("\n".join(out))
print("DONE: ref_deep.txt written, lines=%d" % len(out))
