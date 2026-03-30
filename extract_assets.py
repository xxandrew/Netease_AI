from pptx import Presentation
import os

prs = Presentation(r"C:/Users/xxn4472/Desktop/test1/MC_UXD-AI0323.pptx")
master = prs.slide_masters[0]
out_dir = r"C:\Netease_Ai\03_document_processing\ppt\outputs\assets"
os.makedirs(out_dir, exist_ok=True)

# 提取 master 背景图
for i, sh in enumerate(master.shapes):
    if sh.shape_type == 13:  # PICTURE
        img = sh.image
        ext = img.ext
        path = os.path.join(out_dir, "bg_master.%s" % ext)
        with open(path, "wb") as f:
            f.write(img.blob)
        print("Saved master bg: %s (%.2f x %.2f inches)" % (
            path, sh.width/914400, sh.height/914400))

# 提取 slide 1 中的所有图片
for idx, slide in enumerate(prs.slides):
    for sh in slide.shapes:
        if sh.shape_type == 13:
            img = sh.image
            ext = img.ext
            path = os.path.join(out_dir, "slide%d_%s.%s" % (idx+1, sh.name, ext))
            with open(path, "wb") as f:
                f.write(img.blob)
            print("Saved: %s" % path)
        elif sh.shape_type == 6:
            for ch in sh.shapes:
                if ch.shape_type == 13:
                    img = ch.image
                    ext = img.ext
                    path = os.path.join(out_dir, "slide%d_%s.%s" % (idx+1, ch.name, ext))
                    with open(path, "wb") as f:
                        f.write(img.blob)
                    print("Saved nested: %s" % path)

print("DONE")
