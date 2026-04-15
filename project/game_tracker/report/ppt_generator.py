"""
游戏分析 PPT 生成器
包含：封面、执行摘要、趋势图、区域对比、榜单变动、竞品对比
"""
import io
import os
from datetime import datetime, timedelta

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import matplotlib.font_manager as fm
import numpy as np

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE

from database.db import get_conn

# ── 字体 & 颜色主题 ──────────────────────────────────────────────
BG_DARK   = RGBColor(0x0D, 0x11, 0x17)
BG_CARD   = RGBColor(0x16, 0x1B, 0x22)
ACCENT    = RGBColor(0x58, 0xA6, 0xFF)
GREEN     = RGBColor(0x3F, 0xB9, 0x50)
ORANGE    = RGBColor(0xF0, 0x88, 0x3E)
RED       = RGBColor(0xF8, 0x51, 0x49)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
MUTED     = RGBColor(0x8B, 0x94, 0x9E)
PURPLE    = RGBColor(0xBC, 0x8C, 0xFF)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


# ── matplotlib 中文字体支持 ──────────────────────────────────────
def _setup_matplotlib():
    """尝试设置中文字体，失败则用默认"""
    for font_name in ["Microsoft YaHei", "SimHei", "PingFang SC", "WenQuanYi Micro Hei"]:
        try:
            fm.findfont(fm.FontProperties(family=font_name), fallback_to_default=False)
            plt.rcParams["font.family"] = font_name
            return
        except Exception:
            continue
    plt.rcParams["font.family"] = "DejaVu Sans"

_setup_matplotlib()


# ── 辅助函数 ────────────────────────────────────────────────────

def _add_text(tf, text, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    tf.text = ""
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color

def _add_paragraph(tf, text, size=14, bold=False, color=WHITE, align=PP_ALIGN.LEFT, space_before=0):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return p

def _fill_shape(shape, color: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color

def _add_box(slide, left, top, width, height, bg_color=BG_CARD, radius=False):
    box = slide.shapes.add_shape(1, left, top, width, height)  # MSO_SHAPE_TYPE.RECTANGLE=1
    _fill_shape(box, bg_color)
    box.line.color.rgb = RGBColor(0x30, 0x36, 0x3D)
    box.line.width = Pt(0.5)
    return box

def _img_to_stream(fig):
    buf = io.BytesIO()
    fig.savefig(buf, format="png", bbox_inches="tight",
                facecolor="#0D1117", edgecolor="none", dpi=150)
    buf.seek(0)
    plt.close(fig)
    return buf


# ── 数据查询 ────────────────────────────────────────────────────

def _query_game_data(name: str, days: int = 30):
    """查询指定游戏近 N 天的数据"""
    conn = get_conn()
    c = conn.cursor()
    since = (datetime.now() - timedelta(days=days)).isoformat()

    # Steam 趋势
    steam_trend = c.execute("""
        SELECT fetched_at, current_players, rank, region_code
        FROM steam_games
        WHERE name=? AND fetched_at>=?
        ORDER BY fetched_at ASC
    """, (name, since)).fetchall()

    # 各区最新排名
    region_ranks = c.execute("""
        SELECT region_code, region, rank FROM steam_games s1
        WHERE name=?
          AND fetched_at=(SELECT MAX(fetched_at) FROM steam_games s2 WHERE s2.region_code=s1.region_code)
        ORDER BY rank ASC
    """, (name,)).fetchall()

    # 手游各区排名历史
    mobile_trend = c.execute("""
        SELECT fetched_at, rank, region_code, platform
        FROM mobile_games
        WHERE name=? AND fetched_at>=?
        ORDER BY fetched_at ASC
    """, (name, since)).fetchall()

    # 竞品（同时期全球 Top 20，排除自身）
    competitors = c.execute("""
        SELECT name, current_players, rank FROM steam_games
        WHERE region_code='global' AND name!=?
          AND fetched_at=(SELECT MAX(fetched_at) FROM steam_games WHERE region_code='global')
        ORDER BY rank ASC LIMIT 10
    """, (name,)).fetchall()

    # 榜单变动历史（进入/退出）
    snapshots = c.execute("""
        SELECT DISTINCT fetched_at FROM steam_games
        WHERE region_code='global' ORDER BY fetched_at ASC
    """).fetchall()

    conn.close()
    return {
        "steam_trend": [dict(r) for r in steam_trend],
        "region_ranks": [dict(r) for r in region_ranks],
        "mobile_trend": [dict(r) for r in mobile_trend],
        "competitors": [dict(r) for r in competitors],
        "snapshots": [r[0] for r in snapshots],
    }


def _compute_summary(name, data):
    """计算关键摘要指标"""
    global_trend = [d for d in data["steam_trend"] if d["region_code"] == "global"]
    if global_trend:
        latest = global_trend[-1]
        prev   = global_trend[0] if len(global_trend) > 1 else latest
        peak   = max(d["current_players"] or 0 for d in global_trend)
        avg    = int(sum(d["current_players"] or 0 for d in global_trend) / len(global_trend))
        change = (latest["current_players"] or 0) - (prev["current_players"] or 0)
        best_rank = min((d["rank"] or 999) for d in global_trend if d["rank"])
    else:
        latest = {}; peak = avg = change = 0; best_rank = "-"

    mobile_latest = {}
    for d in reversed(data["mobile_trend"]):
        key = f"{d['region_code']}_{d['platform']}"
        if key not in mobile_latest:
            mobile_latest[key] = d

    return {
        "current_players": latest.get("current_players", 0) or 0,
        "current_rank": latest.get("rank", "-"),
        "peak_players": peak,
        "avg_players": avg,
        "player_change": change,
        "best_rank": best_rank,
        "region_count": len(data["region_ranks"]),
        "mobile_platforms": len(mobile_latest),
    }


# ── 图表生成 ────────────────────────────────────────────────────

def _chart_steam_trend(name, global_trend):
    """Steam 全球在线人数趋势折线图"""
    if not global_trend:
        return None
    dates = [d["fetched_at"][:10] for d in global_trend]
    players = [d["current_players"] or 0 for d in global_trend]

    # 去重日期（取每天最后一条）
    day_map = {}
    for d in global_trend:
        day_map[d["fetched_at"][:10]] = d["current_players"] or 0
    dates = list(day_map.keys())
    players = list(day_map.values())

    fig, ax = plt.subplots(figsize=(10, 3.5))
    fig.patch.set_facecolor("#0D1117")
    ax.set_facecolor("#161B22")

    x = range(len(dates))
    ax.fill_between(x, players, alpha=0.2, color="#58A6FF")
    ax.plot(x, players, color="#58A6FF", linewidth=2.5, marker="o", markersize=4)

    ax.set_xticks(range(0, len(dates), max(1, len(dates)//8)))
    ax.set_xticklabels([dates[i] for i in range(0, len(dates), max(1, len(dates)//8))],
                       color="#8B949E", fontsize=8, rotation=30)
    ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda v, _: f"{int(v/1000)}K" if v >= 1000 else str(int(v))))
    ax.tick_params(colors="#8B949E", labelsize=9)
    for spine in ax.spines.values():
        spine.set_edgecolor("#30363D")
    ax.grid(axis="y", color="#21262D", linewidth=0.7)
    ax.set_title(f"{name} — Steam 全球在线人数趋势", color="white", fontsize=12, pad=10)

    return _img_to_stream(fig)


def _chart_region_rank(data_dict):
    """各区域当前排名柱状图"""
    ranks = data_dict["region_ranks"]
    if not ranks:
        return None

    labels = [r["region"] or r["region_code"] for r in ranks]
    values = [r["rank"] or 0 for r in ranks]
    colors = ["#58A6FF","#3FB950","#BC8CFF","#F0883E","#F85149"]

    fig, ax = plt.subplots(figsize=(8, 3.5))
    fig.patch.set_facecolor("#0D1117")
    ax.set_facecolor("#161B22")

    bars = ax.barh(labels, values, color=colors[:len(labels)], height=0.5)
    ax.invert_yaxis()
    ax.invert_xaxis()  # 排名越小越好，越靠右
    for bar, val in zip(bars, values):
        ax.text(bar.get_width() + 0.5, bar.get_y() + bar.get_height()/2,
                f"#{val}", va="center", color="white", fontsize=10, fontweight="bold")

    ax.tick_params(colors="#8B949E", labelsize=9)
    for spine in ax.spines.values():
        spine.set_edgecolor("#30363D")
    ax.set_xlabel("排名（数字越小=排名越高）", color="#8B949E", fontsize=9)
    ax.set_title("各区域当前排名", color="white", fontsize=12, pad=10)
    ax.grid(axis="x", color="#21262D", linewidth=0.7)

    return _img_to_stream(fig)


def _chart_competitors(name, competitors, current_players):
    """竞品横向对比条形图"""
    items = list(competitors[:8])
    labels = [c["name"][:15] for c in items]
    values = [c["current_players"] or 0 for c in items]

    # 插入目标游戏
    target_val = current_players
    colors = []
    final_labels, final_values = [], []
    inserted = False
    for lbl, val in zip(labels, values):
        if not inserted and target_val >= val:
            final_labels.append(name[:15])
            final_values.append(target_val)
            colors.append("#F0883E")
            inserted = True
        final_labels.append(lbl)
        final_values.append(val)
        colors.append("#58A6FF")
    if not inserted:
        final_labels.append(name[:15])
        final_values.append(target_val)
        colors.append("#F0883E")

    fig, ax = plt.subplots(figsize=(10, 3.8))
    fig.patch.set_facecolor("#0D1117")
    ax.set_facecolor("#161B22")

    y_pos = range(len(final_labels))
    bars = ax.barh(list(y_pos), final_values, color=colors, height=0.6)
    ax.set_yticks(list(y_pos))
    ax.set_yticklabels(final_labels, color="white", fontsize=9)
    ax.invert_yaxis()

    max_val = max(final_values) if final_values else 1
    for bar, val in zip(bars, final_values):
        ax.text(min(val + max_val*0.01, max_val*0.98), bar.get_y() + bar.get_height()/2,
                f"{val:,}", va="center", color="white", fontsize=8)

    ax.xaxis.set_major_formatter(plt.FuncFormatter(lambda v,_: f"{int(v/1000)}K" if v>=1000 else str(int(v))))
    ax.tick_params(colors="#8B949E", labelsize=8)
    for spine in ax.spines.values():
        spine.set_edgecolor("#30363D")
    ax.grid(axis="x", color="#21262D", linewidth=0.7)
    ax.set_title("竞品在线人数横向对比（Steam 全球）", color="white", fontsize=12, pad=10)

    return _img_to_stream(fig)


def _chart_rank_timeline(global_trend):
    """排名变化折线图（排名越小越好，Y轴反转）"""
    ranked = [d for d in global_trend if d.get("rank")]
    if not ranked:
        return None

    day_map = {}
    for d in ranked:
        day_map[d["fetched_at"][:10]] = d["rank"]
    dates = list(day_map.keys())
    ranks = list(day_map.values())

    fig, ax = plt.subplots(figsize=(10, 3))
    fig.patch.set_facecolor("#0D1117")
    ax.set_facecolor("#161B22")

    x = range(len(dates))
    ax.plot(x, ranks, color="#3FB950", linewidth=2, marker="s", markersize=4)
    ax.invert_yaxis()

    ax.set_xticks(range(0, len(dates), max(1, len(dates)//8)))
    ax.set_xticklabels([dates[i] for i in range(0, len(dates), max(1, len(dates)//8))],
                       color="#8B949E", fontsize=8, rotation=30)
    ax.yaxis.set_major_formatter(plt.FuncFormatter(lambda v,_: f"#{int(v)}"))
    ax.tick_params(colors="#8B949E", labelsize=9)
    for spine in ax.spines.values():
        spine.set_edgecolor("#30363D")
    ax.grid(color="#21262D", linewidth=0.7)
    ax.set_title("Steam 全球排名变化趋势", color="white", fontsize=12, pad=10)

    return _img_to_stream(fig)


# ── PPT 构建 ────────────────────────────────────────────────────

def _slide_cover(prs, name, period_label, generated_at):
    """Slide 1: 封面"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = BG_DARK

    # 顶部装饰条
    bar = slide.shapes.add_shape(1, 0, 0, SLIDE_W, Inches(0.08))
    _fill_shape(bar, ACCENT); bar.line.fill.background()

    # 游戏名
    txb = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11), Inches(1.6))
    tf = txb.text_frame; tf.word_wrap = True
    _add_text(tf, name, size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    # 副标题
    txb2 = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(11), Inches(0.7))
    tf2 = txb2.text_frame
    _add_text(tf2, "游戏数据表现分析报告", size=22, color=ACCENT, align=PP_ALIGN.CENTER)

    # 周期 & 生成时间
    txb3 = slide.shapes.add_textbox(Inches(1), Inches(4.4), Inches(11), Inches(0.5))
    tf3 = txb3.text_frame
    _add_text(tf3, f"分析周期：{period_label}　｜　生成时间：{generated_at}", size=14, color=MUTED, align=PP_ALIGN.CENTER)

    # 底部装饰条
    bot = slide.shapes.add_shape(1, 0, SLIDE_H - Inches(0.08), SLIDE_W, Inches(0.08))
    _fill_shape(bot, ACCENT); bot.line.fill.background()

    # 角标
    logo = slide.shapes.add_textbox(Inches(0.3), SLIDE_H - Inches(0.5), Inches(4), Inches(0.35))
    _add_text(logo.text_frame, "🎮 游戏数据追踪看板", size=10, color=MUTED)


def _slide_summary(prs, name, summary):
    """Slide 2: 执行摘要 - 关键指标卡"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = BG_DARK

    # 标题
    _add_slide_title(slide, "📊 执行摘要", "关键数据指标一览")

    # 6 个指标卡
    metrics = [
        ("当前在线人数", f"{summary['current_players']:,}", ACCENT, "Steam 全球实时"),
        ("历史峰值", f"{summary['peak_players']:,}", GREEN, "分析周期内"),
        ("平均在线", f"{summary['avg_players']:,}", PURPLE, "分析周期均值"),
        ("当前排名", f"#{summary['current_rank'] or '-'}", ORANGE, "Steam 全球榜"),
        ("最高排名", f"#{summary['best_rank']}", GREEN, "分析周期内"),
        ("覆盖区域", f"{summary['region_count']} 个", ACCENT, "Steam 各区"),
    ]
    card_w = Inches(3.8)
    card_h = Inches(1.55)
    cols = 3
    start_x = Inches(0.5)
    start_y = Inches(1.6)
    gap_x = Inches(0.35)
    gap_y = Inches(0.3)

    for i, (label, value, color, sub) in enumerate(metrics):
        col = i % cols
        row = i // cols
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)

        box = _add_box(slide, x, y, card_w, card_h)

        # 左侧色条
        accent_bar = slide.shapes.add_shape(1, x, y, Inches(0.06), card_h)
        _fill_shape(accent_bar, color); accent_bar.line.fill.background()

        lbl_box = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.15), card_w - Inches(0.2), Inches(0.35))
        _add_text(lbl_box.text_frame, label, size=11, color=MUTED)

        val_box = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.45), card_w - Inches(0.2), Inches(0.7))
        _add_text(val_box.text_frame, value, size=28, bold=True, color=color)

        sub_box = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(1.1), card_w - Inches(0.2), Inches(0.3))
        _add_text(sub_box.text_frame, sub, size=10, color=MUTED)

    # 变动指示
    chg = summary["player_change"]
    chg_color = GREEN if chg >= 0 else RED
    chg_str = f"{'▲' if chg >= 0 else '▼'} {abs(chg):,}（vs 周期首日）"
    chg_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(12), Inches(0.5))
    _add_text(chg_box.text_frame, f"在线人数变动：{chg_str}", size=13, color=chg_color)


def _slide_trend(prs, name, img_players, img_rank):
    """Slide 3: 趋势折线图"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = BG_DARK
    _add_slide_title(slide, "📈 数据趋势", "在线人数 & 排名变化")

    if img_players:
        slide.shapes.add_picture(img_players, Inches(0.4), Inches(1.4), Inches(12.5), Inches(3.2))
    if img_rank:
        slide.shapes.add_picture(img_rank, Inches(0.4), Inches(4.7), Inches(12.5), Inches(2.5))


def _slide_region(prs, name, img_region, region_ranks):
    """Slide 4: 区域对比"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = BG_DARK
    _add_slide_title(slide, "🌍 区域表现", "各地区 Steam 排名对比")

    if img_region:
        slide.shapes.add_picture(img_region, Inches(0.4), Inches(1.4), Inches(7), Inches(3.5))

    # 右侧表格
    if region_ranks:
        table_x = Inches(7.8)
        table_y = Inches(1.4)
        headers = ["区域", "排名"]
        rows_data = [[r.get("region") or r["region_code"], f"#{r['rank'] or '-'}"] for r in region_ranks[:8]]

        tbl = slide.shapes.add_table(len(rows_data)+1, 2,
                                      table_x, table_y, Inches(4.8), Inches(0.42*(len(rows_data)+1))).table
        for j, h in enumerate(headers):
            cell = tbl.cell(0, j)
            cell.text = h
            cell.text_frame.paragraphs[0].runs[0].font.size = Pt(11)
            cell.text_frame.paragraphs[0].runs[0].font.bold = True
            cell.text_frame.paragraphs[0].runs[0].font.color.rgb = WHITE
            cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0x21, 0x26, 0x2D)

        for i, row in enumerate(rows_data):
            for j, val in enumerate(row):
                cell = tbl.cell(i+1, j)
                cell.text = val
                cell.text_frame.paragraphs[0].runs[0].font.size = Pt(11)
                color = ACCENT if j == 1 else WHITE
                cell.text_frame.paragraphs[0].runs[0].font.color.rgb = color
                cell.fill.solid()
                cell.fill.fore_color.rgb = BG_CARD if i % 2 == 0 else RGBColor(0x1A, 0x1F, 0x28)


def _slide_diff(prs, name, data):
    """Slide 5: 榜单变动时间线"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = BG_DARK
    _add_slide_title(slide, "📋 榜单变动记录", "分析周期内进榜 / 掉榜 / 排名波动")

    global_trend = [d for d in data["steam_trend"] if d["region_code"] == "global"]

    if not global_trend:
        txb = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11), Inches(1))
        _add_text(txb.text_frame, "暂无足够历史数据", size=18, color=MUTED, align=PP_ALIGN.CENTER)
        return

    # 按时间排列变动事件
    events = []
    prev_rank = None
    for d in global_trend:
        rank = d.get("rank")
        ts = d["fetched_at"][:10]
        if prev_rank is None and rank:
            events.append(("🆕 首次进榜", ts, f"排名 #{rank}", GREEN))
        elif prev_rank and not rank:
            events.append(("👻 掉出榜单", ts, f"上次排名 #{prev_rank}", RED))
        elif rank and prev_rank:
            delta = prev_rank - rank
            if abs(delta) >= 5:
                icon = "⬆️ 排名大涨" if delta > 0 else "⬇️ 排名大跌"
                color = GREEN if delta > 0 else RED
                events.append((icon, ts, f"#{prev_rank} → #{rank}（变动 {delta:+d}）", color))
        prev_rank = rank

    if not events:
        events = [("📊 稳定在榜", global_trend[0]["fetched_at"][:10],
                   f"排名范围 #{min(d['rank'] for d in global_trend if d.get('rank'))} ~ #{max(d['rank'] for d in global_trend if d.get('rank'))}",
                   ACCENT)]

    y_start = Inches(1.6)
    line_h = Inches(0.65)
    for i, (icon, ts, detail, color) in enumerate(events[:9]):
        y = y_start + i * line_h
        # 时间线圆点
        dot = slide.shapes.add_shape(9, Inches(0.5), y + Inches(0.18), Inches(0.28), Inches(0.28))
        _fill_shape(dot, color); dot.line.fill.background()
        # 竖线
        if i < len(events) - 1:
            line = slide.shapes.add_shape(1, Inches(0.62), y + Inches(0.45), Inches(0.04), Inches(0.55))
            _fill_shape(line, MUTED); line.line.fill.background()

        # 文本
        txb = slide.shapes.add_textbox(Inches(1.0), y, Inches(11), Inches(0.55))
        tf = txb.text_frame
        _add_text(tf, f"{icon}  {ts}", size=11, bold=True, color=color)
        _add_paragraph(tf, f"  {detail}", size=10, color=MUTED)


def _slide_competitors(prs, name, img_comp, competitors):
    """Slide 6: 竞品对比"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background; bg.fill.solid(); bg.fill.fore_color.rgb = BG_DARK
    _add_slide_title(slide, "🏆 竞品横向对比", "Steam 全球同期热门游戏对比")

    if img_comp:
        slide.shapes.add_picture(img_comp, Inches(0.4), Inches(1.4), Inches(12.5), Inches(4.0))

    # 底部注释
    note = slide.shapes.add_textbox(Inches(0.5), Inches(6.3), Inches(12), Inches(0.5))
    _add_text(note.text_frame, f"* 橙色为 {name}，蓝色为竞品；数据来源：SteamSpy", size=10, color=MUTED)


def _add_slide_title(slide, title, subtitle=""):
    """通用页面标题"""
    # 顶部色条
    bar = slide.shapes.add_shape(1, 0, 0, SLIDE_W, Inches(0.06))
    _fill_shape(bar, ACCENT); bar.line.fill.background()

    txb = slide.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(12), Inches(0.7))
    tf = txb.text_frame
    _add_text(tf, title, size=24, bold=True, color=WHITE)
    if subtitle:
        _add_paragraph(tf, subtitle, size=13, color=MUTED)

    # 分割线
    sep = slide.shapes.add_shape(1, Inches(0.5), Inches(1.1), Inches(12.3), Inches(0.02))
    _fill_shape(sep, RGBColor(0x30, 0x36, 0x3D)); sep.line.fill.background()


# ── 主入口 ──────────────────────────────────────────────────────

def generate_ppt(name: str, days: int = 30) -> io.BytesIO:
    """
    生成指定游戏的分析 PPT，返回 BytesIO 对象
    name: 游戏名称
    days: 分析天数
    """
    data = _query_game_data(name, days)
    summary = _compute_summary(name, data)
    global_trend = [d for d in data["steam_trend"] if d["region_code"] == "global"]

    now_str = datetime.now().strftime("%Y-%m-%d %H:%M")
    period_label = f"近 {days} 天（{(datetime.now()-timedelta(days=days)).strftime('%Y-%m-%d')} ~ {datetime.now().strftime('%Y-%m-%d')}）"

    # 生成图表
    img_players = _chart_steam_trend(name, global_trend)
    img_rank    = _chart_rank_timeline(global_trend)
    img_region  = _chart_region_rank(data)
    img_comp    = _chart_competitors(name, data["competitors"], summary["current_players"])

    # 构建 PPT
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    _slide_cover(prs, name, period_label, now_str)
    _slide_summary(prs, name, summary)
    _slide_trend(prs, name, img_players, img_rank)
    _slide_region(prs, name, img_region, data["region_ranks"])
    _slide_diff(prs, name, data)
    _slide_competitors(prs, name, img_comp, data["competitors"])

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf
